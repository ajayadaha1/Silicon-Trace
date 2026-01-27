"""
PPTX Parser Module for Silicon Trace
Extracts hardware failure data from PowerPoint presentations

Supports two extraction methods:
1. Direct extraction from native PowerPoint tables/text (fast)
2. OCR fallback for image-based content (slower but handles screenshots)
"""

from typing import List, Dict, Any, Optional
from pathlib import Path
import io
import re
from collections import Counter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
from PIL import Image

# OCR is optional - only import if available
try:
    import easyocr
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("Warning: easyocr not installed. OCR fallback will be disabled.")
    print("To enable OCR: pip install easyocr")

from parser import SerialNumberDetector, clean_value, normalize_column_name, extract_customer_from_filename


class PPTXParser:
    """
    Intelligent PowerPoint parser that extracts asset data using multiple strategies
    """
    
    def __init__(self):
        self.ocr_reader = None  # Lazy load OCR reader
        self.customer_from_filename = None  # Store customer extracted from filename
        self.stats = {
            'slides_processed': 0,
            'tables_extracted': 0,
            'text_extracted': 0,
            'ocr_used': 0,
            'assets_found': 0
        }
    
    def parse_pptx(self, file_path: str, original_filename: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        Parse a PPTX file and extract asset data
        
        Args:
            file_path: Path to PPTX file
            original_filename: Original filename for tracking
            
        Returns:
            List of asset dictionaries compatible with main parser format
        """
        file_path_obj = Path(file_path)
        if not file_path_obj.exists():
            raise FileNotFoundError(f"PPTX file not found: {file_path}")
        
        source_filename = original_filename if original_filename else file_path_obj.name
        
        # Extract customer name from filename
        customer_from_filename = extract_customer_from_filename(source_filename)
        if customer_from_filename:
            print(f"PPTX Parser: Extracted customer '{customer_from_filename}' from filename")
        
        # Load presentation
        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise ValueError(f"Error reading PPTX file: {str(e)}")
        
        print(f"PPTX Parser: Processing {len(prs.slides)} slides from '{source_filename}'")
        
        all_data = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            self.stats['slides_processed'] += 1
            slide_data = []
            
            # Phase 1: Try direct table extraction
            tables_data = self._extract_tables_from_slide(slide, slide_idx)
            if tables_data:
                self.stats['tables_extracted'] += len(tables_data)
                slide_data.extend(tables_data)
                print(f"  Slide {slide_idx}: Extracted {len(tables_data)} native tables")
            
            # Phase 1b: Try text extraction (for bullet points, text boxes)
            if not slide_data:
                text_data = self._extract_text_from_slide(slide, slide_idx)
                if text_data:
                    self.stats['text_extracted'] += 1
                    slide_data.extend(text_data)
                    print(f"  Slide {slide_idx}: Extracted text content")
            
            # Phase 2: OCR fallback for image-heavy slides
            if not slide_data:
                ocr_data = self._extract_via_ocr(slide, slide_idx)
                if ocr_data:
                    self.stats['ocr_used'] += 1
                    slide_data.extend(ocr_data)
                    print(f"  Slide {slide_idx}: Used OCR extraction")
            
            if not slide_data:
                print(f"  Slide {slide_idx}: No data extracted (empty or unsupported content)")
            
            all_data.extend(slide_data)
        
        # Convert to standard asset format
        assets = self._convert_to_asset_format(all_data, source_filename, customer_from_filename)
        self.stats['assets_found'] = len(assets)
        
        # Print summary
        print(f"\nPPTX Parser Summary:")
        print(f"  Slides processed: {self.stats['slides_processed']}")
        print(f"  Native tables: {self.stats['tables_extracted']}")
        print(f"  Text extraction: {self.stats['text_extracted']}")
        print(f"  OCR fallback: {self.stats['ocr_used']}")
        print(f"  Assets found: {self.stats['assets_found']}")
        
        return assets
    
    def _extract_tables_from_slide(self, slide, slide_idx: int) -> List[pd.DataFrame]:
        """Extract native PowerPoint tables from a slide"""
        tables = []
        
        for shape in slide.shapes:
            if shape.has_table:
                try:
                    table = shape.table
                    
                    # Extract table data
                    data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        data.append(row_data)
                    
                    if len(data) > 1:  # Need at least header + 1 data row
                        # Convert to DataFrame
                        df = pd.DataFrame(data[1:], columns=data[0])
                        df['_source_slide'] = slide_idx
                        tables.append(df)
                        
                except Exception as e:
                    print(f"    Warning: Failed to extract table from slide {slide_idx}: {e}")
        
        return tables
    
    def _extract_text_from_slide(self, slide, slide_idx: int) -> List[Dict[str, Any]]:
        """Extract text content from slide (bullet points, text boxes)"""
        text_blocks = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    text_blocks.append(text)
        
        if not text_blocks:
            return []
        
        # Combine all text
        full_text = '\n'.join(text_blocks)
        
        # Look for patterns like "9MP1796P50010 (EX HWA)" or "9MP7222Q50001 (SYSTEM_HANG)"
        # Pattern: Serial number (must contain digits or underscores/dashes) followed by text in parentheses
        # Avoid matching pure words like "Collection (text)"
        pattern = r'([A-Z0-9_\-]{8,})\s*\(([^)]+)\)'
        matches = re.findall(pattern, full_text, re.IGNORECASE)
        
        records = []
        if matches:
            # Found serial numbers with error types in parentheses
            # Filter to ensure it's actually a serial number (has digits or special chars, not just letters)
            for serial_num, error_info in matches:
                serial_clean = serial_num.strip()
                # Require at least one digit OR underscore/dash to be a valid serial
                if re.search(r'\d|_|-', serial_clean):
                    records.append({
                        'serial_number': serial_clean,
                        'error_type': error_info.strip(),
                        'raw_text': f"{serial_clean} ({error_info})",
                        '_source_slide': slide_idx,
                        '_extraction_method': 'text'
                    })
        else:
            # Fallback: Look for serial numbers without parentheses
            serial_pattern = r'\b[A-Z0-9]{8,}(?:[_\-][A-Z0-9]+)*\b'
            serials = re.findall(serial_pattern, full_text)
            
            if serials:
                # Parse text into structured data
                lines = full_text.split('\n')
                
                current_record = {'_source_slide': slide_idx, '_extraction_method': 'text'}
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    # Try to parse key-value pairs (e.g., "Serial: ABC123" or "Status: Closed")
                    if ':' in line:
                        parts = line.split(':', 1)
                        key = parts[0].strip()
                        value = parts[1].strip() if len(parts) > 1 else ''
                        if value:
                            current_record[key] = value
                    else:
                        # Check if line contains a serial number
                        if any(sn in line for sn in serials):
                            current_record['raw_text'] = line
                
                if len(current_record) > 2:  # More than just _source_slide and _extraction_method
                    records.append(current_record)
        
        return records
    
    def _extract_via_ocr(self, slide, slide_idx: int) -> List[Dict[str, Any]]:
        """Extract data using OCR (for image-based content)"""
        if not OCR_AVAILABLE:
            print(f"    Slide {slide_idx}: OCR not available (easyocr not installed)")
            return []
        
        try:
            # Convert slide to image
            image = self._slide_to_image(slide)
            if image is None:
                return []
            
            # Initialize OCR reader if needed (lazy load)
            if self.ocr_reader is None:
                print("    Initializing OCR engine (this may take a moment)...")
                self.ocr_reader = easyocr.Reader(['en'], gpu=False)
            
            # Perform OCR
            results = self.ocr_reader.readtext(image)
            
            # Extract text with confidence filtering
            texts = []
            for (bbox, text, confidence) in results:
                if confidence > 0.5:  # Filter low-confidence results
                    texts.append(text)
            
            if not texts:
                return []
            
            # Combine OCR results
            full_text = ' '.join(texts)
            
            # Parse into structured data (similar to text extraction)
            return self._parse_ocr_text(full_text, slide_idx)
            
        except Exception as e:
            print(f"    Warning: OCR failed for slide {slide_idx}: {e}")
            return []
    
    def _slide_to_image(self, slide) -> Optional[Image.Image]:
        """Convert slide to PIL Image (placeholder - needs implementation)"""
        # Note: python-pptx doesn't directly support slide rendering
        # This would require an external tool like:
        # - LibreOffice in headless mode
        # - unoconv
        # - comtypes (Windows only)
        # For now, we'll skip this and implement if OCR is needed
        return None
    
    def _parse_ocr_text(self, text: str, slide_idx: int) -> List[Dict[str, Any]]:
        """Parse OCR text into structured records"""
        records = []
        
        # Look for serial numbers
        serial_pattern = r'\b[A-Z0-9]{8,}(?:[_\-][A-Z0-9]+)*\b'
        serials = re.findall(serial_pattern, text)
        
        if serials:
            # Create basic record with detected serials
            for sn in serials:
                records.append({
                    'serial_number': sn,
                    'ocr_text': text,
                    '_source_slide': slide_idx,
                    '_extraction_method': 'ocr'
                })
        
        return records
    
    def _convert_to_asset_format(self, data: List, source_filename: str, customer_from_filename: Optional[str] = None) -> List[Dict[str, Any]]:
        """Convert extracted data to standard asset format"""
        assets = []
        
        for item in data:
            if isinstance(item, pd.DataFrame):
                # Handle DataFrame (from tables)
                df_assets = self._dataframe_to_assets(item, source_filename, customer_from_filename)
                assets.extend(df_assets)
            elif isinstance(item, dict):
                # Handle dict (from text/OCR)
                asset = self._dict_to_asset(item, source_filename, customer_from_filename)
                if asset:
                    assets.append(asset)
        
        return assets
    
    def _dataframe_to_assets(self, df: pd.DataFrame, source_filename: str, customer_from_filename: Optional[str] = None) -> List[Dict[str, Any]]:
        """Convert DataFrame to asset records"""
        assets = []
        
        # Debug: Print all columns
        print(f"    Table columns: {list(df.columns)}")
        
        # Detect serial number column
        serial_column = SerialNumberDetector.detect_serial_column(df)
        if not serial_column:
            print(f"    Warning: No serial number column detected in table")
            return assets
        
        print(f"    Detected serial column: '{serial_column}'")
        
        # Detect error, status, and other important columns with better heuristics
        error_column = None
        status_column = None
        component_column = None
        
        # Error type keywords (English and Chinese)
        error_keywords = ['error', 'failure', 'issue', 'symptom', 'problem', 'fail', 
                         'type', 'failure type', '故障', '错误', 'defect']
        
        # Status keywords (English and Chinese)
        status_keywords = ['status', 'state', 'condition', '状态', 'fa status']
        
        # Component keywords
        component_keywords = ['component', 'part', 'child', 'subpart']
        
        for col in df.columns:
            col_lower = col.lower().strip()
            
            # Detect error column
            if not error_column and any(kw in col_lower for kw in error_keywords):
                error_column = col
            
            # Detect status column
            if not status_column and any(kw in col_lower for kw in status_keywords):
                status_column = col
            
            # Detect component column
            if not component_column and any(kw in col_lower for kw in component_keywords):
                component_column = col
        
        # Process each row
        for idx, row in df.iterrows():
            serial_number = str(row[serial_column]).strip()
            
            # Skip invalid serial numbers
            if not serial_number or serial_number.lower() in ['nan', 'none', '', 'null', 'nat', 'n/a', 'na', 'tbd', 'tbc']:
                continue
            
            # Build raw_data - preserve ALL columns
            raw_data = {}
            for col in df.columns:
                value = clean_value(row[col])
                if value is not None:
                    raw_data[col] = value
            
            # Add metadata
            raw_data['_source_slide'] = row.get('_source_slide', 'unknown')
            raw_data['_extraction_method'] = 'native_table'
            
            # Add customer from filename as fallback if no customer column exists
            if customer_from_filename:
                has_customer_column = any('customer' in normalize_column_name(col).lower() for col in df.columns)
                if not has_customer_column:
                    raw_data['Customer'] = customer_from_filename
            
            # Extract error_type with smart handling
            error_type = None
            if error_column and error_column in row:
                error_value = clean_value(row[error_column])
                if error_value and str(error_value).strip():
                    error_type = str(error_value).strip()
            
            # Extract status with smart handling
            status = None
            if status_column and status_column in row:
                status_value = clean_value(row[status_column])
                if status_value and str(status_value).strip():
                    status = str(status_value).strip()
            
            asset = {
                'serial_number': serial_number,
                'error_type': error_type,
                'status': status,
                'source_filename': source_filename,
                'raw_data': raw_data
            }
            
            assets.append(asset)
        
        return assets
    
    def _dict_to_asset(self, data: Dict[str, Any], source_filename: str, customer_from_filename: Optional[str] = None) -> Optional[Dict[str, Any]]:
        """Convert dict to asset record"""
        # Check if serial_number is already in the dict (from improved text extraction)
        serial_number = data.get('serial_number')
        
        if not serial_number:
            # Try to find serial number in dict keys
            for key, value in data.items():
                key_lower = key.lower()
                if 'serial' in key_lower or 'sn' in key_lower or 'cpu' in key_lower:
                    serial_number = str(value).strip()
                    break
        
        if not serial_number:
            # Try to find in raw_text
            if 'raw_text' in data:
                serial_pattern = r'\b[A-Z0-9]{8,}(?:[_\-][A-Z0-9]+)*\b'
                matches = re.findall(serial_pattern, data['raw_text'])
                if matches:
                    serial_number = matches[0]
        
        if not serial_number:
            return None
        
        # Build asset
        raw_data = {k: v for k, v in data.items() if not k.startswith('_')}
        raw_data['_source_slide'] = data.get('_source_slide', 'unknown')
        raw_data['_extraction_method'] = data.get('_extraction_method', 'text')
        
        # Add customer from filename as fallback if no customer data exists
        if customer_from_filename:
            has_customer = any('customer' in normalize_column_name(k).lower() for k in raw_data.keys())
            if not has_customer:
                raw_data['Customer'] = customer_from_filename
        
        asset = {
            'serial_number': serial_number,
            'error_type': None,
            'status': None,
            'source_filename': source_filename,
            'raw_data': raw_data
        }
        
        return asset


# Convenience function for external use
def parse_pptx(file_path: str, original_filename: Optional[str] = None) -> List[Dict[str, Any]]:
    """
    Parse a PPTX file and extract asset data
    
    Args:
        file_path: Path to PPTX file
        original_filename: Original filename for tracking
        
    Returns:
        List of asset dictionaries
    """
    parser = PPTXParser()
    return parser.parse_pptx(file_path, original_filename)
